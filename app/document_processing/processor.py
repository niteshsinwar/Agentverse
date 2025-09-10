"""
Document Processing Module
Handles extraction of text content from various file formats
Uses lightweight libraries for better reliability
"""
import os
import pandas as pd
from pathlib import Path
from typing import Optional
import tempfile
import base64

# Core libraries
from PIL import Image
import fitz  # PyMuPDF for PDF and image processing
import pypdf
from docx import Document
from pptx import Presentation

# OpenAI for intelligent document analysis
try:
    import openai
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False


class DocumentProcessor:
    """Processes various document formats and extracts text content"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._extract_pdf_content,
            '.docx': self._extract_docx_content,
            '.pptx': self._extract_pptx_content,
            '.csv': self._extract_csv_content,
            '.txt': self._extract_text_content,
            '.jpg': self._extract_image_content,
            '.jpeg': self._extract_image_content,
            '.png': self._extract_image_content,
            '.gif': self._extract_image_content
        }
        
        # Initialize OpenAI client if available
        self.openai_client = None
        if OPENAI_AVAILABLE:
            try:
                # Try to get API key from environment
                api_key = os.getenv('OPENAI_API_KEY')
                if api_key:
                    self.openai_client = OpenAI(api_key=api_key)
            except Exception as e:
                print(f"Warning: Could not initialize OpenAI client: {e}")
    
    def process_file(self, file_path: str, file_extension: str = None, use_ai_analysis: bool = True) -> str:
        """Process a file and extract its text content"""
        if file_extension is None:
            file_extension = Path(file_path).suffix.lower()
        
        if file_extension not in self.supported_formats:
            return f"Unsupported file format: {file_extension}"
        
        try:
            processor = self.supported_formats[file_extension]
            basic_content = processor(file_path)
            
            # Enhance with AI analysis if available and requested
            if use_ai_analysis and self.openai_client:
                try:
                    ai_content = self._ai_analyze_content(file_path, file_extension, basic_content)
                    if ai_content:
                        return f"{basic_content}\n\n--- AI Analysis ---\n{ai_content}"
                except Exception as e:
                    print(f"Warning: AI analysis failed: {e}")
            
            return basic_content
        except Exception as e:
            return f"Error processing {file_extension} file: {str(e)}"
    
    def _ai_analyze_content(self, file_path: str, file_extension: str, basic_content: str) -> str:
        """Use OpenAI to analyze document content"""
        if not self.openai_client:
            return ""
        
        try:
            # Handle images with vision model
            if file_extension.lower() in ['.jpg', '.jpeg', '.png', '.gif']:
                return self._ai_analyze_image(file_path)
            
            # Handle text-based documents
            elif len(basic_content.strip()) > 0:
                return self._ai_analyze_text_document(basic_content, file_extension)
                
        except Exception as e:
            print(f"AI analysis error: {e}")
            return ""
    
    def _ai_analyze_image(self, file_path: str) -> str:
        """Analyze image using OpenAI Vision"""
        try:
            # Encode image to base64
            with open(file_path, "rb") as image_file:
                base64_image = base64.b64encode(image_file.read()).decode('utf-8')
            
            response = self.openai_client.chat.completions.create(
                model="gpt-4o",  # or "gpt-4-vision-preview"
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Analyze this image thoroughly. Describe what you see, extract any text content, identify key elements, and provide insights about the purpose or context of this image."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{base64_image}"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=1000
            )
            
            return response.choices[0].message.content
            
        except Exception as e:
            print(f"Image analysis error: {e}")
            return f"Could not analyze image: {str(e)}"
    
    def _ai_analyze_text_document(self, content: str, file_extension: str) -> str:
        """Analyze text document using OpenAI"""
        try:
            file_type = file_extension.upper().lstrip('.')
            
            prompt = f"""Analyze this {file_type} document content and provide:
1. A comprehensive summary
2. Key insights and important information
3. Main topics or themes
4. Any actionable items or important data points
5. Structure and organization analysis

Document content:
{content[:4000]}  # Limit content to avoid token limits
"""
            
            response = self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "user", "content": prompt}
                ],
                max_tokens=800
            )
            
            return response.choices[0].message.content
            
        except Exception as e:
            print(f"Text analysis error: {e}")
            return f"Could not analyze document: {str(e)}"
    
    def _extract_pdf_content(self, file_path: str) -> str:
        """Extract text from PDF using PyMuPDF (more reliable than PyPDF2)"""
        try:
            content = []
            
            # Try PyMuPDF first (handles both text and image-based PDFs)
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    content.append(f"--- Page {page_num + 1} ---\n{text}")
                else:
                    # If no text found, try to extract from images on the page
                    try:
                        pix = page.get_pixmap()
                        # Convert to PIL Image for processing
                        img_data = pix.tobytes("ppm")
                        with tempfile.NamedTemporaryFile(suffix=".ppm", delete=False) as temp_img:
                            temp_img.write(img_data)
                            temp_img.flush()
                            
                            # Use PyMuPDF's basic text recognition if available
                            # Note: This is basic and may not work for all cases
                            content.append(f"--- Page {page_num + 1} (Image-based) ---\n[Image content detected but text extraction limited without OCR]")
                            os.unlink(temp_img.name)
                    except Exception:
                        content.append(f"--- Page {page_num + 1} ---\n[Unable to extract content from this page]")
            
            doc.close()
            return '\n\n'.join(content) if content else "No readable content found in PDF"
            
        except Exception as e:
            # Fallback to pypdf if PyMuPDF fails
            try:
                content = []
                with open(file_path, 'rb') as file:
                    reader = pypdf.PdfReader(file)
                    for page_num, page in enumerate(reader.pages):
                        text = page.extract_text()
                        if text.strip():
                            content.append(f"--- Page {page_num + 1} ---\n{text}")
                
                return '\n\n'.join(content) if content else "No readable content found in PDF"
            except Exception as e2:
                return f"Error extracting PDF content: {str(e2)}"
    
    def _extract_docx_content(self, file_path: str) -> str:
        """Extract text from Word documents"""
        try:
            doc = Document(file_path)
            content = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_content = []
                    for cell in row.cells:
                        row_content.append(cell.text.strip())
                    table_content.append(' | '.join(row_content))
                
                if table_content:
                    content.append('\n--- Table ---\n' + '\n'.join(table_content))
            
            return '\n\n'.join(content) if content else "No readable content found in document"
            
        except Exception as e:
            return f"Error extracting DOCX content: {str(e)}"
    
    def _extract_pptx_content(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            prs = Presentation(file_path)
            content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                
                if len(slide_content) > 1:  # More than just the header
                    content.append('\n'.join(slide_content))
            
            return '\n\n'.join(content) if content else "No readable content found in presentation"
            
        except Exception as e:
            return f"Error extracting PPTX content: {str(e)}"
    
    def _extract_csv_content(self, file_path: str) -> str:
        """Extract and summarize CSV data"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                return "Error: Could not read CSV file with standard encodings"
            
            content = []
            content.append(f"CSV Data Summary:")
            content.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
            content.append(f"Columns: {', '.join(df.columns.tolist())}")
            
            # Show first few rows
            if len(df) > 0:
                content.append("\n--- First 5 rows ---")
                content.append(df.head().to_string(index=False))
            
            # Basic statistics for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                content.append("\n--- Numeric Summary ---")
                content.append(df[numeric_cols].describe().to_string())
            
            return '\n'.join(content)
            
        except Exception as e:
            return f"Error extracting CSV content: {str(e)}"
    
    def _extract_text_content(self, file_path: str) -> str:
        """Extract content from plain text files"""
        try:
            encodings = ['utf-8', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            
            return "Error: Could not read text file with standard encodings"
            
        except Exception as e:
            return f"Error extracting text content: {str(e)}"
    
    def _extract_image_content(self, file_path: str) -> str:
        """Extract text from images using PyMuPDF (limited without full OCR)"""
        try:
            # Open image with PyMuPDF
            doc = fitz.open(file_path)
            page = doc.load_page(0)  # Images have only one "page"
            
            # Get image properties
            pix = page.get_pixmap()
            width, height = pix.width, pix.height
            
            # Basic image analysis
            content = []
            content.append(f"Image Analysis:")
            content.append(f"Dimensions: {width} x {height} pixels")
            content.append(f"Format: {Path(file_path).suffix.upper()}")
            
            # Try to detect if image contains text (basic analysis)
            try:
                # Convert to PIL for basic analysis
                img_data = pix.tobytes("ppm")
                with tempfile.NamedTemporaryFile(suffix=".ppm", delete=False) as temp_img:
                    temp_img.write(img_data)
                    temp_img.flush()
                    
                    # Open with PIL for basic analysis
                    pil_img = Image.open(temp_img.name)
                    
                    # Basic heuristics to detect if image might contain text
                    # (This is very basic - for full OCR, you'd need Tesseract or cloud OCR)
                    if width > 100 and height > 50:  # Reasonable size for text
                        content.append("\n[Image appears to be of suitable size for text content]")
                        content.append("[Note: Full text extraction requires OCR capabilities]")
                        content.append("[To extract text from images, consider using cloud OCR services]")
                    else:
                        content.append("\n[Image may be too small to contain readable text]")
                    
                    os.unlink(temp_img.name)
                    
            except Exception:
                content.append("\n[Unable to analyze image content]")
            
            doc.close()
            return '\n'.join(content)
            
        except Exception as e:
            return f"Error analyzing image: {str(e)}"
